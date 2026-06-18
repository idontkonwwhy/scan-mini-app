import os, json
from pptx import Presentation

PPT_DIR = r"C:\Users\Yux\Desktop\家具配件数据"
OUT_DIR = r"E:\scan-mini-app\scripts"

ppts = [
    ("ppt1", "兴盛发五金产品图册-1（三合一锁扣，板托，三合一轮，滑轨，脚轮，胶链，胶链灯,吊码）(1).pptx"),
    ("ppt2", "兴盛发五金产品图册-2 (吊码，调节家具脚，赌头，反弹器，防水垫，防撞角，挂衣挂，衣托，合页，透气网，家具锁，毛边，螺母，快装杆，拉直器)(1).pptx"),
    ("ppt3", "兴盛发五金产品图册-3 (拉直器，铁衣架，铁鞋架，伸宿衣架，蛇形管，连接杆，暗藏板托，玻璃扣，角码，床角撑，线盒，孔塞，盖子，下水口，拉手，地钉)(1).pptx"),
    ("ppt4", "兴盛发五金产品图册-4（毛条，合页孔塞，酒杯架，隐藏胶链，十字胶，透明橡胶家具脚垫）(1).pptx"),
]

all_results = {}

for ppt_key, filename in ppts:
    filepath = os.path.join(PPT_DIR, filename)
    if not os.path.exists(filepath):
        print(f"NOT FOUND: {filepath}")
        continue

    print(f"\n==== {ppt_key}: {filename} ====")
    prs = Presentation(filepath)
    slide_texts = {}

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1  # 1-indexed
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        full_text = " | ".join(texts)
        slide_texts[f"slide{slide_num:02d}"] = full_text
        if full_text:
            print(f"  slide{slide_num:02d}: {full_text[:120]}{'...' if len(full_text)>120 else ''}")
        else:
            print(f"  slide{slide_num:02d}: (空)")

    all_results[ppt_key] = slide_texts

# Save extracted data
out_path = os.path.join(OUT_DIR, "ppt-slide-texts.json")
with open(out_path, 'w', encoding='utf-8') as f:
    json.dump(all_results, f, ensure_ascii=False, indent=2)

print(f"\n\nExtracted {sum(len(v) for v in all_results.values())} slides to {out_path}")
